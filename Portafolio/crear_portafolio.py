import os
from tkinter import messagebox
import docx
from docx.shared import Pt, Cm
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
import datetime
from manipular_json import manipularJson
from docx2pdf import convert
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
class crearPortafolio():
    LISTA_CARPETAS = ['PORTADA', 
                    'DESCRIPCIÓN DEL CURSO',
                    'PRESENTACIÓN GENERAL DEL ESTUDIANTE',
                    'ACTIVIDADES O ASIGNACIONES', 
                    'MATERIAL DIDACTICO DEL CURSO',
                    'CONCLUSIÓN']
    def __init__(self, ruta_predefinda, nombre_carpeta) -> None:
        self.ruta_carpeta = ruta_predefinda
        self.nombre_carpeta = nombre_carpeta
        self.crear_carpetas(self.ruta_carpeta, self.nombre_carpeta)

    def crear_carpetas(self, ruta_carpeta, nombreCarpeta):
        #Crear carpeta principal del portafolio
        try:
            if not self.ruta_carpeta:
                messagebox.showwarning('Advertencia', 'Se debe selecciona una ruta donde crear el portafolio.')
                return
            else:
                ruta_completa = os.path.join(ruta_carpeta, nombreCarpeta)
                os.makedirs(ruta_completa)
        except FileExistsError:
            messagebox.showwarning('Advertencia', 'El portafolio que estás intentando crear ya existe, prueba con otro nombre.')
            return
        #Carear subcarpetas del portafolio
        try:
            if os.path.exists(ruta_completa):
                for subcarpetas in self.LISTA_CARPETAS:
                    ruta_subcarpetas = os.path.join(ruta_completa, subcarpetas)
                    os.makedirs(ruta_subcarpetas)
            else:
                messagebox.showwarning('Advertencia', 'La carpeta principal del portafolio no existe, intentalo de nuevo.')
                return
        except Exception as h:
            print(h)
        
        try:
            if os.path.exists(os.path.join(ruta_completa, self.LISTA_CARPETAS[3])):
                cargar_datos = manipularJson()
                credenciales_usuario = cargar_datos.cargar_credenciales()
                for nombre, condicion in credenciales_usuario['Subcarpeta_actividades'].items():
                    ruta_subcarpetas = os.path.join(ruta_completa, self.LISTA_CARPETAS[3])
                    if condicion:
                        os.makedirs(os.path.join(ruta_subcarpetas, nombre.upper()))
        except Exception as e:
            print(e)

    def mover_archivos(self, carpeta_origen, carpetaDestino):
        if os.path.exists(carpeta_origen):
            pass

    def crear_portada(self):
        self.doc = docx.Document()

        cargar_datos = manipularJson()
        credenciales_usuario = cargar_datos.cargar_credenciales()
        facultad_seleccionada = credenciales_usuario['Credenciales']['Facultad']
       # Crear una tabla con una fila y tres celdas
        table = self.doc.add_table(rows=1, cols=3)
        table.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        table.autofit = False
        # Obtener la primera celda (izquierda) y añadir una imagen
        cell_izquierda = table.cell(0, 0)
        cell_izquierda.width = Cm(2.3)
        run_izquierda = cell_izquierda.paragraphs[0].add_run()
        imagen_izquierda = run_izquierda.add_picture(r'.\Portafolio\imagenes\logo_utp.png', width=Cm(2.3))

        # Obtener la segunda celda (centro) y agregar el texto
        cell_centro = table.cell(0, 1)
        parrafo_centro = cell_centro.add_paragraph('Universidad Tecnológica de Panamá')
        run = parrafo_centro.runs[0]
        font = run.font
        font.name = 'Arial'
        font.size = Pt(13)
        parrafo_centro.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        
        parrafo_centro = cell_centro.add_paragraph(credenciales_usuario['Credenciales']['Facultad'])
        run = parrafo_centro.runs[0]
        font = run.font
        font.name = 'Arial'
        font.size = Pt(13)
        parrafo_centro.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        if facultad_seleccionada == 'Facultad de Ingeniería de Sistemas Computacionales':  
            parrafo_centro = cell_centro.add_paragraph(credenciales_usuario['Credenciales']['Departamento'])
            run = parrafo_centro.runs[0]
            font = run.font
            font.name = 'Arial'
            font.size = Pt(13)
            parrafo_centro.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        parrafo_centro = cell_centro.add_paragraph(credenciales_usuario['Credenciales']['Carrera'])
        run = parrafo_centro.runs[0]
        font = run.font
        font.name = 'Arial'
        font.size = Pt(13)
        parrafo_centro.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        #Espacio
        parrafo_centro = cell_centro.add_paragraph()
        parrafo_centro.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        parrafo_centro = cell_centro.add_paragraph(credenciales_usuario['Credenciales']['Materia'])
        run = parrafo_centro.runs[0]
        font = run.font
        font.name = 'Arial'
        font.size = Pt(13)
        parrafo_centro.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        #Espacio
        parrafo_centro = cell_centro.add_paragraph()
        parrafo_centro.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        parrafo_centro = cell_centro.add_paragraph('Estudiante:')
        run = parrafo_centro.runs[0]
        font = run.font
        font.name = 'Arial'
        font.size = Pt(13)
        parrafo_centro.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        
        parrafo_centro = cell_centro.add_paragraph(
            credenciales_usuario['Credenciales']['Apellido']+', '+
            credenciales_usuario['Credenciales']['Nombre']+' '+
            credenciales_usuario['Credenciales']['Cedula'])
        run = parrafo_centro.runs[0]
        font = run.font
        font.name = 'Arial'
        font.size = Pt(13)
        parrafo_centro.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        #Espacio
        for _ in range (2):
            parrafo_centro = cell_centro.add_paragraph()
            parrafo_centro.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
            
        parrafo_centro = cell_centro.add_paragraph('Profesor:')
        run = parrafo_centro.runs[0]
        font = run.font
        font.name = 'Arial'
        font.size = Pt(13)
        parrafo_centro.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        
        parrafo_centro = cell_centro.add_paragraph(credenciales_usuario['Credenciales']['Profesor'])
        run = parrafo_centro.runs[0]
        font = run.font
        font.name = 'Arial'
        font.size = Pt(13)
        parrafo_centro.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        for _ in range (4):
            parrafo_centro = cell_centro.add_paragraph()
            parrafo_centro.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        parrafo_centro = cell_centro.add_paragraph(credenciales_usuario['Credenciales']['Grupo'])
        run = parrafo_centro.runs[0]
        font = run.font
        font.name = 'Arial'
        font.size = Pt(13)
        parrafo_centro.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        parrafo_centro = cell_centro.add_paragraph(self.obtener_semestre())
        run = parrafo_centro.runs[0]
        font = run.font
        font.name = 'Arial'
        font.size = Pt(13)
        parrafo_centro.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        #Ancho de la celda del medio
        cell_centro.width = Cm(13)
        # Obtener la tercera celda (derecha) y añadir una imagen
        cell_derecha = table.cell(0, 2)
        cell_derecha.width = Cm(2.4)
        run_derecha = cell_derecha.paragraphs[0].add_run()
        if facultad_seleccionada == 'Facultad de Ingeniería de Sistemas Computacionales':
            logo_facultad = r'.\Portafolio\imagenes\logo_fisc.png'
        elif facultad_seleccionada == 'Facultad de Ingeniería Civil':
            logo_facultad = r'.\Portafolio\imagenes\logo_civil.jpg'
        elif facultad_seleccionada == 'Facultad de Ingeniería Eléctrica':
            logo_facultad =r'.\Portafolio\imagenes\logo_electrica.jpg'
        elif facultad_seleccionada == 'Facultad de Ingeniería Industrial':
            logo_facultad = r'.\Portafolio\imagenes\logo_industrial.png'
        elif facultad_seleccionada == 'Facultad de Ingeniería Mecánica':
            logo_facultad = r'.\Portafolio\imagenes\logo_mecanica.png'
        elif facultad_seleccionada == 'Facultad de Ciencias y Tecnología':
            logo_facultad = r'.\Portafolio\imagenes\logo_ciencias_tecnologia.png'
        imagen_derecha = run_derecha.add_picture(logo_facultad, width=Cm(2.2))
        # Ajustar los márgenes de las celdas para que las imágenes estén a la izquierda y a la derecha
        cell_izquierda.paragraphs[0].alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        cell_derecha.paragraphs[0].alignment = WD_PARAGRAPH_ALIGNMENT.RIGHT
        
        ruta = os.path.join(self.ruta_carpeta, self.nombre_carpeta)
        ruta_guardado = str(os.path.join(ruta, 'PORTADA')+'\Portada.docx')
        try:
            self.doc.save(ruta_guardado)
            convert(ruta_guardado)
            os.remove(ruta_guardado)
        except Exception as e:
            messagebox.showwarning('Advertencia', e)
  
    def obtener_semestre(self):
        fecha_actual = datetime.datetime.now()
        if 8 <= fecha_actual.month <=12:
            return str('Semestre II, '+str(fecha_actual.year))
        elif 1 <= fecha_actual.month <= 3:
            return str('Verano, '+str(fecha_actual.year))
        else:
            return  str('Semestre I, '+str(fecha_actual.year))
    
    def crear_presentacion_personal(self):
        cargar_datos = manipularJson()
        credenciales_usuario = cargar_datos.cargar_credenciales()
        ppt = pptx.Presentation()
        slide_layout = ppt.slide_layouts[1]
        slide = ppt.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        print(content.text_frame.paragraphs[0].level)
        title.text = str(credenciales_usuario['Credenciales']['Nombre']+' '+credenciales_usuario['Credenciales']['Apellido'])
        content.text = str(credenciales_usuario['Credenciales']['Intereses'])
        img_path = str(credenciales_usuario['Credenciales']['Ruta Foto de Perfil'])
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.name = 'Arial'
        content.text_frame.paragraphs[0].font.size = Pt(16)
        content.text_frame.paragraphs[0].font.name = 'Arial'
        left, top, width, height = Inches(0.5), Inches(2), Inches(3.5), Inches(3.5)
        slide.shapes.add_picture(img_path, left, top, width, height)
        left_content = Inches(4.2)
        top_content = top  
        width_content = Inches(5) 
        height_content = height
        content.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.JUSTIFY
        content.left = left_content
        content.top = top_content
        content.width = width_content
        content.height = height_content
        # Configurar el formato del párrafo para eliminar las viñetas
        ppt.save('Presentacion Personal.pptx')

if __name__ == '__main__':
    test = crearPortafolio(r'C:\Users\HP Envy\Downloads', 'Portafolio de Ejemplo')
    #test.crear_portada()
    test.crear_presentacion_personal()